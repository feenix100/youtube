#!/usr/bin/env python3
"""
wildfire_report.py
Download NASA FIRMS active fire detections (CSV) for a region & day range,
export to Excel, generate a detailed map-style scatter (PNG), and create a PowerPoint summary.
Requires free api map key from Nasa: https://firms.modaps.eosdis.nasa.gov/api/map_key/

FIRMS = Fire Information for Resource Management System

Coordinate data: https://download.geonames.org/export/dump/

Examples:
  # Read MAP_KEY from ./MAP_KEY file, world, last 3 days
  python wildfire_report_firms_min.py --days 3

  # Western US bounding box (west,south,east,north), last 2 days
  python wildfire_report_firms_min.py --bbox "-130,30,-100,50" --days 2

  create a virtual env then install requirements to use script:

  python -m venv NameOfVenv

  pip install -r requirements.txt 

"""
import argparse, os, sys, io
from pathlib import Path
from textwrap import dedent

import requests
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib as mpl
import matplotlib.dates as mdates
from mpl_toolkits.axes_grid1.inset_locator import inset_axes
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.util import Inches, Pt

# ---- FIRMS API ----
FIRMS_AREA_API = "https://firms.modaps.eosdis.nasa.gov/api/area/csv"
# Common FIRMS sources:
#   MODIS_NRT, MODIS_SP,
#   VIIRS_SNPP_NRT, VIIRS_SNPP_SP, VIIRS_NOAA20_NRT, VIIRS_NOAA20_SP, VIIRS_NOAA21_NRT, LANDSAT_NRT (US/CA)
DEFAULT_SOURCE = "VIIRS_SNPP_NRT"

# ---- URL builder ----
def build_url(mapkey: str, source: str, area: str, days: int, date: str | None) -> str:
    # /api/area/csv/{MAP_KEY}/{SOURCE}/{AREA_COORDINATES}/{DAY_RANGE}[/{DATE}]
    base = f"{FIRMS_AREA_API}/{mapkey}/{source}/{area}/{days}"
    return f"{base}/{date}" if date else base

# ---- MAP_KEY loader ----
def load_map_key(cli_key: str | None, key_file: str | None) -> str:
    """Resolve MAP_KEY from CLI arg, environment, or a local file (default 'MAP_KEY')."""
    if cli_key and cli_key.strip():
        return cli_key.strip()
    env = os.getenv("FIRMS_MAP_KEY") or os.getenv("MAP_KEY")
    if env and env.strip():
        return env.strip()
    path = Path(key_file or "MAP_KEY")
    if path.is_file():
        key = path.read_text(encoding="utf-8").strip()
        if key:
            return key
    raise RuntimeError(
        "MAP_KEY not found. Provide one by either:\n"
        "  • --mapkey YOUR_MAP_KEY\n"
        "  • export FIRMS_MAP_KEY=YOUR_MAP_KEY (or MAP_KEY)\n"
        "  • create a file named 'MAP_KEY' with your key (one line)"
    )

# ---- Fetch FIRMS CSV ----
def fetch_fires(mapkey: str, source: str, bbox: str | None, days: int, date: str | None) -> pd.DataFrame:
    if days < 1 or days > 10:
        raise ValueError("--days must be between 1 and 10 (FIRMS limit).")
    area = "world" if not bbox else bbox  # west,south,east,north OR 'world'

    url = build_url(mapkey, source, area, days, date)
    resp = requests.get(url, timeout=60)
    resp.raise_for_status()

    text = resp.text
    # Guard: HTML/JSON indicates error (bad key/params)
    if text.lstrip().startswith("<") or text.strip().startswith("{"):
        raise RuntimeError(
            "FIRMS returned a non-CSV response. Check your MAP_KEY/source/params.\n"
            f"URL: {url}\nFirst 200 chars:\n{text[:200]}"
        )

    try:
        df = pd.read_csv(io.StringIO(text))
    except Exception as e:
        snippet = text.splitlines()[:10]
        raise RuntimeError(f"CSV parse error: {e}\nFirst lines:\n" + "\n".join(snippet)) from e

    if df.empty:
        raise RuntimeError("No detections returned. Try increasing --days or changing --source/--bbox.")
    return df

# ---- Time normalization ----
def normalize_times(df: pd.DataFrame) -> pd.DataFrame:
    """
    Combine FIRMS acq_date (YYYY-MM-DD) & acq_time (hhmm UTC) into naive-UTC 'time_utc' for Excel.
    """
    if "acq_date" in df.columns and "acq_time" in df.columns:
        timestr = df["acq_date"].astype(str) + " " + df["acq_time"].astype(str).str.zfill(4)
        ts = pd.to_datetime(timestr, format="%Y-%m-%d %H%M", utc=True)
        df["time_utc"] = ts.dt.tz_localize(None)  # keep UTC clock time, strip tzinfo for Excel
    return df

# ---- Choose an intensity column for sizing/color ----
def choose_intensity_column(df: pd.DataFrame) -> str:
    for col in ["frp", "brightness", "bright_ti4", "bright_ti5"]:
        if col in df.columns:
            return col
    if "confidence" in df.columns:
        return "confidence"
    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    if not num_cols:
        raise RuntimeError("No numeric columns found to use as intensity.")
    return num_cols[0]

# ---- Save Excel ----
def save_to_excel(
    df: pd.DataFrame,
    path: str,
    *,
    also_txt: bool = False,
    txt_path: str | None = None,
    txt_sep: str = "\t"
):
    """
    Save Excel workbook (Fires + Summary). Optionally also write a tab-delimited TXT.

    Args:
        path: Excel output path (e.g., 'fires.xlsx')
        also_txt: If True, write a TXT export of the 'Fires' DataFrame
        txt_path: Optional explicit TXT path; defaults to same stem with .txt
        txt_sep: Delimiter for TXT (default: tab)
    """
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)

    with pd.ExcelWriter(path, engine="xlsxwriter", datetime_format="yyyy-mm-dd hh:mm") as wr:
        # Raw data
        df.to_excel(wr, sheet_name="Fires", index=False)

        # Summary
        intensity = choose_intensity_column(df)
        summary = {
            "detections": len(df),
            "intensity_metric": intensity,
            "min_intensity": float(df[intensity].min()) if intensity in df else np.nan,
            "max_intensity": float(df[intensity].max()) if intensity in df else np.nan,
            "mean_intensity": float(df[intensity].mean()) if intensity in df else np.nan,
            "start_utc": df["time_utc"].min() if "time_utc" in df else pd.NaT,
            "end_utc": df["time_utc"].max() if "time_utc" in df else pd.NaT,
        }
        pd.DataFrame([summary]).to_excel(wr, sheet_name="Summary", index=False)

        # Autosize columns
        for name, _df in [("Fires", df), ("Summary", pd.DataFrame([summary]))]:
            ws = wr.sheets[name]
            for i, col in enumerate(_df.columns):
                width = min(max(len(str(col)), _df[col].astype(str).map(len).max()), 60)
                ws.set_column(i, i, width + 2)

    # Optional TXT export (Fires only)
    if also_txt:
        out_txt = txt_path or str(Path(path).with_suffix(".txt"))
        os.makedirs(os.path.dirname(out_txt) or ".", exist_ok=True)
        df.to_csv(out_txt, sep=txt_sep, index=False, encoding="utf-8", lineterminator="\n")
        return {"excel": path, "txt": out_txt}

    return {"excel": path}

# ---- Detailed map-like scatter ----
def plot_map(df: pd.DataFrame, png: str, title: str):
    """
    Presence-only map:
      - Plots unique wildfire detection coordinates as uniform points
      - No colorbar, no size legend—just 'where fires were recorded'
    """

    os.makedirs(os.path.dirname(png) or ".", exist_ok=True)
    if not {"latitude", "longitude"}.issubset(df.columns):
        raise RuntimeError("CSV missing 'latitude'/'longitude' columns.")

    # Keep unique, valid coordinates only
    coords = (
        df[["longitude", "latitude"]]
        .dropna()
        .astype({"longitude": float, "latitude": float})
        .drop_duplicates()
    )

    if coords.empty:
        raise RuntimeError("No valid coordinates to plot.")

    lon = coords["longitude"].to_numpy()
    lat = coords["latitude"].to_numpy()

    fig, ax = plt.subplots(figsize=(11, 5.8))
    ax.scatter(lon, lat, s=10, alpha=0.8, marker="o", linewidths=0)  # uniform points

    # Titles / labels
    ax.set_title(f"{title} — {len(coords)} unique locations")
    ax.set_xlabel("Longitude")
    ax.set_ylabel("Latitude")
    ax.grid(True, linestyle="--", alpha=0.35)

    # Fit bounds with gentle padding; handle single-point extents
    x_min, x_max = float(np.min(lon)), float(np.max(lon))
    y_min, y_max = float(np.min(lat)), float(np.max(lat))
    if x_min == x_max:
        pad_x = 0.25
        x_min, x_max = x_min - pad_x, x_max + pad_x
    else:
        pad_x = (x_max - x_min) * 0.03
        x_min, x_max = x_min - pad_x, x_max + pad_x
    if y_min == y_max:
        pad_y = 0.25
        y_min, y_max = y_min - pad_y, y_max + pad_y
    else:
        pad_y = (y_max - y_min) * 0.03
        y_min, y_max = y_min - pad_y, y_max + pad_y

    ax.set_xlim(x_min, x_max)
    ax.set_ylim(y_min, y_max)

    fig.tight_layout()
    fig.savefig(png, dpi=240)
    plt.close(fig)

# ---- Summaries / PPT ----
def summarize(df: pd.DataFrame) -> dict:
    intensity = choose_intensity_column(df)
    mx_idx = df[intensity].astype(float).idxmax()
    top = df.loc[mx_idx]
    # Prefer any descriptive text field; fallback to coordinates
    label = top.get("place") or top.get("location") or ""
    return {
        "count": int(len(df)),
        "intensity_col": intensity,
        "max_intensity": float(top[intensity]),
        "top_loc": (float(top.get("longitude", np.nan)), float(top.get("latitude", np.nan))),
        "top_label": label,
        "start": df["time_utc"].min() if "time_utc" in df else None,
        "end": df["time_utc"].max() if "time_utc" in df else None,
    }

def make_ppt(df: pd.DataFrame, png_path: str, pptx_path: str, region_label: str, days: int, source: str):
    stats = summarize(df)
    prs = Presentation()

    # Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = f"Wildfire Detections — {region_label}"
    s1.placeholders[1].text = f"{source} • Last {days} day(s) • {stats['count']} detections"

    # Map + highlights
    s2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    s2.shapes.title.text = "Active Fire Map & Highlights"
    s2.shapes.add_picture(png_path, Inches(0.5), Inches(1.4), height=Inches(4.2))

    tx = s2.shapes.add_textbox(Inches(6.3), Inches(1.4), Inches(3.2), Inches(4.4))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Highlights"
    p.font.size = Pt(20)

    loc_text = (
        stats["top_label"]
        if stats["top_label"]
        else f"lon {stats['top_loc'][0]:.2f}, lat {stats['top_loc'][1]:.2f}"
    )

    bullets = []
    if stats.get("start") is not None and stats.get("end") is not None:
        bullets.append(
            f"Period: {stats['start'].strftime('%Y-%m-%d %H:%M')} → {stats['end'].strftime('%Y-%m-%d %H:%M')} UTC"
        )
    bullets.extend([
        f"Detections: {stats['count']}",
        f"Intensity metric: {stats['intensity_col']}",
        f"Max intensity: {stats['max_intensity']:.2f} at {loc_text}",
    ])
    for b in bullets:
        pb = tf.add_paragraph()
        pb.text = b
        pb.level = 1

    # Methods
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Data & Methodology"
    s3.placeholders[1].text = dedent(f"""
        • Source: NASA FIRMS Active Fires (CSV via area API)
        • Query: {source}, last {days} day(s), region: {region_label}
        • Fields: lat/lon, acquisition time, confidence, intensity (FRP/brightness)
        • Processing: pandas → Excel; matplotlib scatter (size ~ intensity; color ~ recency)
        • Notes: Times are UTC and timezone-stripped for Excel compatibility
    """).strip()

    prs.save(pptx_path)

# ---- CLI / main ----
def main():
    ap = argparse.ArgumentParser(description="FIRMS wildfire report → Excel, PNG map, PowerPoint.")
    ap.add_argument("--mapkey", help="Your FIRMS MAP_KEY (optional; overrides file/env).")
    ap.add_argument("--mapkey-file", default="MAP_KEY", help="Path to file holding MAP_KEY (default: ./MAP_KEY)")
    ap.add_argument("--source", default=DEFAULT_SOURCE, help=f"FIRMS source (default: {DEFAULT_SOURCE})")
    ap.add_argument("--bbox", default=None, help='Bounding box "west,south,east,north" or omit for world')
    ap.add_argument("--days", type=int, default=3, help="Number of days (1..10)")
    ap.add_argument("--date", default=None, help="Optional start date YYYY-MM-DD (if omitted, most recent days)")
    ap.add_argument("--excel", default="fires.xlsx", help="Excel output path")
    ap.add_argument("--png", default="fires_map.png", help="Map image output path")
    ap.add_argument("--pptx", default="fires_report.pptx", help="PowerPoint output path")
    args = ap.parse_args()

    try:
        mapkey = load_map_key(args.mapkey, args.mapkey_file)
        region_label = "World" if not args.bbox else args.bbox

        df = fetch_fires(mapkey, args.source, args.bbox, args.days, args.date)
        df = normalize_times(df)

        #save_to_excel(df, args.excel)
        save_to_excel(df, "fires.xlsx", also_txt=True, txt_path="fires.txt")
        plot_map(df, args.png, f"Wildfire Detections — {region_label} — {args.source} — {args.days}d")
        make_ppt(df, args.png, args.pptx, region_label, args.days, args.source)
        print("Done!")
        print(f" Excel: {args.excel}\n PNG: {args.png}\n PPTX: {args.pptx}")
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
